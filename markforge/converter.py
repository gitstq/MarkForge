"""
MarkForge Converter Engine
Markdown to multiple formats conversion engine
"""

import os
import re
import json
import shutil
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any, Union
from dataclasses import dataclass
from enum import Enum


class OutputFormat(Enum):
    """Supported output formats"""
    HTML = "html"
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TXT = "txt"
    RTF = "rtf"


@dataclass
class ConversionOptions:
    """Conversion options"""
    output_format: OutputFormat = OutputFormat.HTML
    theme: str = "default"
    highlight_code: bool = True
    include_toc: bool = True
    custom_css: Optional[str] = None
    template: Optional[str] = None
    title: Optional[str] = None
    author: Optional[str] = None
    page_size: str = "A4"
    margin: str = "20mm"


class MarkForgeConverter:
    """
    MarkForge Converter - Convert Markdown to multiple formats
    
    Features:
    - HTML with syntax highlighting and themes
    - PDF with custom styling
    - DOCX with templates
    - PPTX from markdown slides (--- separator)
    - Plain text extraction
    - Batch conversion support
    """
    
    # Built-in themes
    THEMES = {
        "default": {
            "primary_color": "#2563eb",
            "background": "#ffffff",
            "text_color": "#1f2937",
            "code_bg": "#f3f4f6",
            "border_color": "#e5e7eb"
        },
        "dark": {
            "primary_color": "#3b82f6",
            "background": "#1f2937",
            "text_color": "#f9fafb",
            "code_bg": "#374151",
            "border_color": "#4b5563"
        },
        "github": {
            "primary_color": "#0969da",
            "background": "#ffffff",
            "text_color": "#24292f",
            "code_bg": "#f6f8fa",
            "border_color": "#d0d7de"
        },
        "notion": {
            "primary_color": "#0066cc",
            "background": "#ffffff",
            "text_color": "#37352f",
            "code_bg": "#f7f6f3",
            "border_color": "#e3e2de"
        },
        "solarized": {
            "primary_color": "#268bd2",
            "background": "#fdf6e3",
            "text_color": "#657b83",
            "code_bg": "#eee8d5",
            "border_color": "#93a1a1"
        }
    }
    
    def __init__(self, options: Optional[ConversionOptions] = None):
        """Initialize converter with options"""
        self.options = options or ConversionOptions()
        self._markdown_content = ""
        self._metadata = {}
        
    def load_file(self, filepath: Union[str, Path]) -> "MarkForgeConverter":
        """Load markdown from file"""
        filepath = Path(filepath)
        if not filepath.exists():
            raise FileNotFoundError(f"File not found: {filepath}")
        
        with open(filepath, "r", encoding="utf-8") as f:
            content = f.read()
        
        # Extract front matter (YAML-like metadata)
        self._metadata = self._extract_front_matter(content)
        self._markdown_content = self._remove_front_matter(content)
        return self
    
    def load_string(self, content: str) -> "MarkForgeConverter":
        """Load markdown from string"""
        self._metadata = self._extract_front_matter(content)
        self._markdown_content = self._remove_front_matter(content)
        return self
    
    def _extract_front_matter(self, content: str) -> Dict[str, str]:
        """Extract YAML-like front matter from markdown"""
        metadata = {}
        if content.startswith("---"):
            parts = content.split("---", 2)
            if len(parts) >= 3:
                front_matter = parts[1].strip()
                for line in front_matter.split("\n"):
                    if ":" in line:
                        key, value = line.split(":", 1)
                        metadata[key.strip()] = value.strip()
        return metadata
    
    def _remove_front_matter(self, content: str) -> str:
        """Remove front matter from markdown content"""
        if content.startswith("---"):
            parts = content.split("---", 2)
            if len(parts) >= 3:
                return parts[2].strip()
        return content
    
    def convert(self, output_format: Optional[OutputFormat] = None) -> str:
        """Convert markdown to specified format"""
        fmt = output_format or self.options.output_format
        
        converters = {
            OutputFormat.HTML: self._to_html,
            OutputFormat.PDF: self._to_pdf,
            OutputFormat.DOCX: self._to_docx,
            OutputFormat.PPTX: self._to_pptx,
            OutputFormat.TXT: self._to_txt,
            OutputFormat.RTF: self._to_rtf
        }
        
        if fmt not in converters:
            raise ValueError(f"Unsupported format: {fmt}")
        
        return converters[fmt]()
    
    def convert_to_file(self, output_path: Union[str, Path], 
                        output_format: Optional[OutputFormat] = None) -> Path:
        """Convert and save to file"""
        output_path = Path(output_path)
        content = self.convert(output_format)
        
        # Ensure parent directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Handle binary formats
        fmt = output_format or self.options.output_format
        if fmt in [OutputFormat.PDF, OutputFormat.DOCX, OutputFormat.PPTX]:
            mode = "wb"
        else:
            mode = "w"
        
        with open(output_path, mode if fmt in [OutputFormat.PDF, OutputFormat.DOCX, OutputFormat.PPTX] else "w", 
                  encoding="utf-8" if mode == "w" else None) as f:
            f.write(content)
        
        return output_path
    
    def _to_html(self) -> str:
        """Convert markdown to HTML with styling"""
        theme = self.THEMES.get(self.options.theme, self.THEMES["default"])
        
        # Parse markdown to HTML (simple implementation)
        html_body = self._parse_markdown_to_html()
        
        # Generate complete HTML document
        html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{self._metadata.get("title", "Document")}</title>
    <style>
        :root {{
            --primary-color: {theme["primary_color"]};
            --background: {theme["background"]};
            --text-color: {theme["text_color"]};
            --code-bg: {theme["code_bg"]};
            --border-color: {theme["border_color"]};
        }}
        
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
            line-height: 1.6;
            color: var(--text-color);
            background: var(--background);
            max-width: 900px;
            margin: 0 auto;
            padding: 40px 20px;
        }}
        
        h1, h2, h3, h4, h5, h6 {{
            margin-top: 1.5em;
            margin-bottom: 0.5em;
            font-weight: 600;
            line-height: 1.25;
        }}
        
        h1 {{ font-size: 2em; border-bottom: 2px solid var(--primary-color); padding-bottom: 0.3em; }}
        h2 {{ font-size: 1.5em; border-bottom: 1px solid var(--border-color); padding-bottom: 0.3em; }}
        h3 {{ font-size: 1.25em; }}
        
        p {{ margin: 1em 0; }}
        
        a {{
            color: var(--primary-color);
            text-decoration: none;
        }}
        
        a:hover {{ text-decoration: underline; }}
        
        code {{
            background: var(--code-bg);
            padding: 0.2em 0.4em;
            border-radius: 3px;
            font-family: "SFMono-Regular", Consolas, "Liberation Mono", Menlo, monospace;
            font-size: 0.9em;
        }}
        
        pre {{
            background: var(--code-bg);
            padding: 16px;
            border-radius: 6px;
            overflow-x: auto;
            margin: 1em 0;
        }}
        
        pre code {{
            background: none;
            padding: 0;
        }}
        
        blockquote {{
            border-left: 4px solid var(--primary-color);
            padding-left: 1em;
            margin: 1em 0;
            color: #6b7280;
        }}
        
        ul, ol {{
            padding-left: 2em;
            margin: 1em 0;
        }}
        
        li {{ margin: 0.5em 0; }}
        
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 1em 0;
        }}
        
        th, td {{
            border: 1px solid var(--border-color);
            padding: 8px 12px;
            text-align: left;
        }}
        
        th {{
            background: var(--code-bg);
            font-weight: 600;
        }}
        
        tr:nth-child(even) {{ background: var(--code-bg); }}
        
        hr {{
            border: none;
            border-top: 1px solid var(--border-color);
            margin: 2em 0;
        }}
        
        img {{
            max-width: 100%;
            height: auto;
            border-radius: 6px;
        }}
        
        .toc {{
            background: var(--code-bg);
            padding: 20px;
            border-radius: 6px;
            margin-bottom: 2em;
        }}
        
        .toc h2 {{
            margin-top: 0;
            border-bottom: none;
        }}
        
        .toc ul {{ list-style: none; padding-left: 1em; }}
        
        @media print {{
            body {{ max-width: none; padding: 0; }}
        }}
    </style>
</head>
<body>
    {self._generate_toc() if self.options.include_toc else ""}
    {html_body}
</body>
</html>'''
        return html
    
    def _parse_markdown_to_html(self) -> str:
        """Parse markdown content to HTML (simple implementation)"""
        html = self._markdown_content
        
        # Escape HTML entities first
        html = html.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        
        # Code blocks (must be first to prevent other parsing inside)
        html = re.sub(
            r'```(\w*)\n(.*?)```',
            lambda m: f'<pre><code class="language-{m.group(1)}">{m.group(2)}</code></pre>',
            html,
            flags=re.DOTALL
        )
        
        # Inline code
        html = re.sub(r'`([^`]+)`', r'<code>\1</code>', html)
        
        # Headers
        html = re.sub(r'^######\s+(.+)$', r'<h6>\1</h6>', html, flags=re.MULTILINE)
        html = re.sub(r'^#####\s+(.+)$', r'<h5>\1</h5>', html, flags=re.MULTILINE)
        html = re.sub(r'^####\s+(.+)$', r'<h4>\1</h4>', html, flags=re.MULTILINE)
        html = re.sub(r'^###\s+(.+)$', r'<h3>\1</h3>', html, flags=re.MULTILINE)
        html = re.sub(r'^##\s+(.+)$', r'<h2>\1</h2>', html, flags=re.MULTILINE)
        html = re.sub(r'^#\s+(.+)$', r'<h1>\1</h1>', html, flags=re.MULTILINE)
        
        # Bold and italic
        html = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', html)
        html = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html)
        html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html)
        html = re.sub(r'___(.+?)___', r'<strong><em>\1</em></strong>', html)
        html = re.sub(r'__(.+?)__', r'<strong>\1</strong>', html)
        html = re.sub(r'_(.+?)_', r'<em>\1</em>', html)
        
        # Links
        html = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<a href="\2">\1</a>', html)
        
        # Images
        html = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', r'<img src="\2" alt="\1">', html)
        
        # Blockquotes
        html = re.sub(r'^>\s+(.+)$', r'<blockquote>\1</blockquote>', html, flags=re.MULTILINE)
        
        # Horizontal rules
        html = re.sub(r'^---$', r'<hr>', html, flags=re.MULTILINE)
        html = re.sub(r'^\*\*\*$', r'<hr>', html, flags=re.MULTILINE)
        
        # Unordered lists
        html = re.sub(r'^[\*\-]\s+(.+)$', r'<li>\1</li>', html, flags=re.MULTILINE)
        
        # Ordered lists
        html = re.sub(r'^\d+\.\s+(.+)$', r'<li>\1</li>', html, flags=re.MULTILINE)
        
        # Tables (simple)
        def parse_table(match):
            lines = match.group(0).strip().split('\n')
            if len(lines) < 2:
                return match.group(0)
            
            header = lines[0].split('|')[1:-1]
            rows = [l.split('|')[1:-1] for l in lines[2:] if l.strip() and not re.match(r'^[\|\-\s]+$', l)]
            
            table_html = '<table>\n<thead>\n<tr>'
            for cell in header:
                table_html += f'<th>{cell.strip()}</th>'
            table_html += '</tr>\n</thead>\n<tbody>\n'
            
            for row in rows:
                table_html += '<tr>'
                for cell in row:
                    table_html += f'<td>{cell.strip()}</td>'
                table_html += '</tr>\n'
            
            table_html += '</tbody>\n</table>'
            return table_html
        
        html = re.sub(r'(\|.+\|[\n\r]+)+', parse_table, html)
        
        # Paragraphs
        paragraphs = html.split('\n\n')
        html = '\n'.join(
            p if p.strip().startswith('<') else f'<p>{p}</p>'
            for p in paragraphs if p.strip()
        )
        
        return html
    
    def _generate_toc(self) -> str:
        """Generate table of contents"""
        headings = re.findall(r'^(#{1,6})\s+(.+)$', self._markdown_content, flags=re.MULTILINE)
        
        if not headings:
            return ""
        
        toc_items = []
        for level, title in headings:
            level_num = len(level)
            indent = "  " * (level_num - 1)
            anchor = re.sub(r'[^\w\-]', '', title.lower().replace(' ', '-'))
            toc_items.append(f'{indent}<li><a href="#{anchor}">{title}</a></li>')
        
        return f'''<nav class="toc">
    <h2>📑 目录</h2>
    <ul>
        {chr(10).join(toc_items)}
    </ul>
</nav>'''
    
    def _to_pdf(self) -> bytes:
        """Convert markdown to PDF (requires wkhtmltopdf or weasyprint)"""
        # Generate HTML first
        html = self._to_html()
        
        # Try to use weasyprint if available
        try:
            from weasyprint import HTML, CSS
            from weasyprint.text.fonts import FontConfiguration
            
            font_config = FontConfiguration()
            doc = HTML(string=html).render(font_config=font_config)
            
            # Return PDF bytes
            import io
            pdf_buffer = io.BytesIO()
            doc.write_pdf(pdf_buffer)
            return pdf_buffer.getvalue()
            
        except ImportError:
            # Fallback: return HTML with PDF content-type hint
            # In production, this would require wkhtmltopdf
            return html.encode('utf-8')
    
    def _to_docx(self) -> bytes:
        """Convert markdown to DOCX (requires python-docx)"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.style import WD_STYLE_TYPE
            
            doc = Document()
            
            # Set document properties
            if self.options.title or self._metadata.get("title"):
                doc.core_properties.title = self.options.title or self._metadata.get("title", "")
            if self.options.author or self._metadata.get("author"):
                doc.core_properties.author = self.options.author or self._metadata.get("author", "")
            
            # Parse and add content
            lines = self._markdown_content.split('\n')
            i = 0
            
            while i < len(lines):
                line = lines[i]
                
                # Headers
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    text = line.lstrip('#').strip()
                    
                    if level == 1:
                        doc.add_heading(text, level=0)
                    else:
                        doc.add_heading(text, level=level)
                
                # Code blocks
                elif line.startswith('```'):
                    lang = line[3:].strip()
                    code_lines = []
                    i += 1
                    while i < len(lines) and not lines[i].startswith('```'):
                        code_lines.append(lines[i])
                        i += 1
                    
                    code_para = doc.add_paragraph()
                    code_run = code_para.add_run('\n'.join(code_lines))
                    code_run.font.name = 'Courier New'
                    code_run.font.size = Pt(9)
                
                # Lists
                elif line.startswith(('- ', '* ', '1. ')):
                    text = line.lstrip('-*1234567890. ').strip()
                    doc.add_paragraph(text, style='List Bullet' if line[0] in '-*' else 'List Number')
                
                # Blockquotes
                elif line.startswith('> '):
                    text = line[2:].strip()
                    para = doc.add_paragraph()
                    para.paragraph_format.left_indent = Inches(0.5)
                    run = para.add_run(text)
                    run.italic = True
                
                # Regular paragraphs
                elif line.strip():
                    # Handle inline formatting
                    text = line
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # Bold (simplified)
                    text = re.sub(r'\*(.+?)\*', r'\1', text)  # Italic (simplified)
                    text = re.sub(r'`(.+?)`', r'\1', text)  # Code (simplified)
                    doc.add_paragraph(text)
                
                i += 1
            
            # Save to bytes
            import io
            buffer = io.BytesIO()
            doc.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            # Fallback: return plain text
            return self._to_txt().encode('utf-8')
    
    def _to_pptx(self) -> bytes:
        """Convert markdown to PPTX (slides separated by ---)"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            
            # Split content by --- for slides
            slides_content = re.split(r'^---$', self._markdown_content, flags=re.MULTILINE)
            
            for slide_content in slides_content:
                slide_content = slide_content.strip()
                if not slide_content:
                    continue
                
                # Create slide
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                # Parse title and content
                lines = slide_content.split('\n')
                title = ""
                content_lines = []
                
                for line in lines:
                    if line.startswith('# ') and not title:
                        title = line[2:].strip()
                    elif not line.startswith('#'):
                        content_lines.append(line.lstrip('-*').strip())
                
                # Set title
                if slide.shapes.title:
                    slide.shapes.title.text = title or "Slide"
                
                # Add content
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Content placeholder
                        tf = shape.text_frame
                        for i, line in enumerate(content_lines[:5]):  # Limit to 5 bullets
                            if i == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = line
                            p.level = 0
                        break
            
            # Save to bytes
            import io
            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            return self._to_txt().encode('utf-8')
    
    def _to_txt(self) -> str:
        """Convert markdown to plain text"""
        text = self._markdown_content
        
        # Remove code block markers
        text = re.sub(r'```\w*\n?', '', text)
        
        # Remove inline code markers
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        # Remove formatting
        text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        text = re.sub(r'_(.+?)_', r'\1', text)
        
        # Convert links
        text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'\1 (\2)', text)
        
        # Remove images
        text = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', r'[Image: \1]', text)
        
        # Clean up headers
        text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
        
        # Clean up blockquotes
        text = re.sub(r'^>\s+', '', text, flags=re.MULTILINE)
        
        # Clean up lists
        text = re.sub(r'^[\*\-]\s+', '• ', text, flags=re.MULTILINE)
        text = re.sub(r'^\d+\.\s+', '', text, flags=re.MULTILINE)
        
        return text.strip()
    
    def _to_rtf(self) -> str:
        """Convert markdown to RTF"""
        # RTF header
        rtf = r'''{{\rtf1\ansi\deff0
{{\fonttbl{{\f0 Times New Roman;}}{{\f1 Courier New;}}}}
{{\colortbl;\red0\green0\blue0;\red0\green0\blue255;}}
\viewkind4\uc1\pard\f0\fs24
'''
        
        # Convert content
        lines = self._markdown_content.split('\n')
        
        for line in lines:
            # Headers
            if line.startswith('#'):
                level = len(line) - len(line.lstrip('#'))
                text = line.lstrip('#').strip()
                size = 48 - (level * 6)  # Larger for higher levels
                rtf += f'\\fs{size} \\b {self._escape_rtf(text)}\\b0\\fs24\\par\n'
            
            # Code blocks
            elif line.startswith('```'):
                continue  # Skip code block markers
            
            # Regular text
            else:
                text = self._escape_rtf(line)
                text = re.sub(r'\*\*(.+?)\*\*', r'\\b \1\\b0 ', text)
                text = re.sub(r'\*(.+?)\*', r'\\i \1\\i0 ', text)
                rtf += f'{text}\\par\n'
        
        rtf += '}'
        return rtf
    
    def _escape_rtf(self, text: str) -> str:
        """Escape special characters for RTF"""
        return text.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
    
    @classmethod
    def batch_convert(cls, input_dir: Union[str, Path], output_dir: Union[str, Path],
                      output_format: OutputFormat = OutputFormat.HTML,
                      options: Optional[ConversionOptions] = None) -> List[Path]:
        """
        Batch convert all markdown files in a directory
        
        Args:
            input_dir: Directory containing markdown files
            output_dir: Directory for output files
            output_format: Output format
            options: Conversion options
            
        Returns:
            List of output file paths
        """
        input_dir = Path(input_dir)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        output_files = []
        
        for md_file in input_dir.glob('*.md'):
            converter = cls(options)
            converter.load_file(md_file)
            
            # Determine output extension
            ext_map = {
                OutputFormat.HTML: '.html',
                OutputFormat.PDF: '.pdf',
                OutputFormat.DOCX: '.docx',
                OutputFormat.PPTX: '.pptx',
                OutputFormat.TXT: '.txt',
                OutputFormat.RTF: '.rtf'
            }
            
            output_file = output_dir / (md_file.stem + ext_map.get(output_format, '.html'))
            converter.convert_to_file(output_file, output_format)
            output_files.append(output_file)
        
        return output_files
